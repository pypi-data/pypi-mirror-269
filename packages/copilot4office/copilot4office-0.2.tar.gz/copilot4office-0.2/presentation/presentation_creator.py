# you can run this test from project root folder as: python -m unittest tests.test_presentation_creator
# or set pythonpath and then run this test from inside tests folder as: python -m unittest test_presentation_creator

import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt

class PresentationCreator:
    """Class to handle creating PowerPoint presentations from JSON data."""
    
    def __init__(self, file_path):
        self.file_path = file_path
        self.presentation = Presentation()
    
    def load_content_from_json(self):
        """Loads content from a JSON file and returns the data."""
        try:
            with open(self.file_path, 'r') as file:
                return json.load(file)
        except FileNotFoundError:
            print(f"Error: The file '{self.file_path}' does not exist.")
            return None
        except json.JSONDecodeError:
            print(f"Error: The file '{self.file_path}' contains invalid JSON.")
            return None
    
    def validate_content_data(self, content_data):
        """Validates the loaded JSON data for necessary keys and structure."""
        if not content_data:
            print("No content data provided.")
            return False
        if 'slides' not in content_data:
            print("Error: No 'slides' key found in the content data.")
            return False
        return all(self.validate_slide(slide) for slide in content_data['slides'])
    
    def validate_slide(self, slide):
        """Validates individual slide data for required keys and image paths."""
        if 'title' not in slide or 'content' not in slide:
            print("Warning: Missing 'title' or 'content' in some slides.")
            return False
        if 'images' in slide and not all('path' in image for image in slide['images']):
            print("Warning: Missing 'path' for an image.")
            return False
        return True
    
    def create_presentation(self, content_data, file_name):
        """Creates a PowerPoint presentation based on the validated JSON content."""
        for index, slide_data in enumerate(content_data['slides'], 1):
            slide = self.add_slide(index)
            self.handle_slide_content(slide, slide_data)
            if 'images' in slide_data:
                img_layout = slide_data.get('imagelayout', [1, 1])
                self.arrange_images_on_slide(slide, slide_data['images'], img_layout)
            self.add_slide_number(slide, index)
        self.presentation.save(file_name)
    
    def add_slide(self, slide_index):
        """Adds a slide to the presentation based on the layout index."""
        layout_index = 0 if slide_index == 1 else 1
        slide_layout = self.presentation.slide_layouts[layout_index]
        return self.presentation.slides.add_slide(slide_layout)
    
    def handle_slide_content(self, slide, slide_data):
        """Populates a slide with title and content."""
        slide.shapes.title.text = slide_data['title']
        slide.placeholders[1].text = slide_data['content']
    
    def arrange_images_on_slide(self, slide, images, layout):
        """Arranges images on a slide based on specified layout and image data."""
        num_rows, num_cols = layout
        for index, image_info in enumerate(images):
            self.arrange_image(slide, image_info, index, num_cols, num_rows)
    
    def arrange_image(self, slide, image_info, index, num_cols, num_rows):
        """Places an individual image on the slide at calculated position."""
        col_offset = index % num_cols
        row_offset = index // num_cols
        content_width = self.presentation.slide_width / num_cols
        row_height = self.presentation.slide_height / num_rows - Inches(4)  # Adjust height
        left = col_offset * content_width
        top = row_offset * row_height + Inches(2)  # Adjust top position
        self.add_image_with_caption(slide, image_info['path'], image_info.get('caption', ""), left, top, content_width / num_cols, row_height / num_rows)
    
    def add_image_with_caption(self, slide, path, caption, left, top, width, height):
        """Adds an image and optionally a caption to a slide."""
        if os.path.exists(path):
            picture = slide.shapes.add_picture(path, left, top, width, height)
            self.add_caption(slide, caption, left, top + height, width)
        else:
            print(f"Warning: Image file '{path}' not found or path is missing.")
    
    def add_caption(self, slide, caption, left, top, width):
        """Adds a caption below an image."""
        caption_box = slide.shapes.add_textbox(left, top, width, Inches(0.5))
        caption_box.text = caption
        caption_box.text_frame.paragraphs[0].font.bold = True
        caption_box.text_frame.paragraphs[0].font.size = Pt(10)
    
    def add_slide_number(self, slide, number):
        """Adds a slide number to the bottom right of a slide."""
        footer_box = slide.shapes.add_textbox(self.presentation.slide_width - Inches(2), self.presentation.slide_height - Inches(0.8), Inches(2), Inches(0.5))
        footer_box.text = str(number)
        footer_box.text_frame.paragraphs[0].alignment = 2  # Right alignment

# Example Usage
# creator = PresentationCreator('path_to_your_json_file.json')
# content = creator.load_content_from_json()
# if creator.validate_content_data(content):
#     creator.create_presentation(content, "output_presentation.pptx")
